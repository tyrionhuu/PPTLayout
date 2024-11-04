from unittest.mock import MagicMock

import pytest
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Length

from pptlayout.extractors.ppt_extractor import (
    PowerPointShapeExtractor,
    SlideShapeExtractor,
)


# Mocking the shape extractor factory function
def mock_shape_extractor_factory(shape):
    mock_extractor = MagicMock()
    mock_extractor.extract_shape.return_value = {
        "shape_id": shape.shape_id,
        "type": shape.shape_type,
    }
    return mock_extractor


# Patch the shape_extractor_factory with our mock
@pytest.fixture(autouse=True)
def patch_shape_extractor_factory(monkeypatch):
    monkeypatch.setattr(
        "pptlayout.extractors.factories.shape_extractor_factory",
        mock_shape_extractor_factory,
    )


# Test SlideShapeExtractor
def test_slide_shape_extractor():
    # Create a mock slide
    mock_slide = MagicMock(spec=Slide)
    mock_slide.slide_id = 1
    mock_slide.name = "Test Slide"
    mock_slide.shapes = [
        MagicMock(shape_id=1, shape_type="AUTO_SHAPE"),
        MagicMock(shape_id=2, shape_type="TEXT_BOX"),
    ]

    slide_extractor = SlideShapeExtractor(mock_slide)

    # Extract slide data
    slide_data = slide_extractor.extract_slide()

    # Verify slide metadata
    assert slide_data["slide_id"] == 1
    assert slide_data["slide_name"] == "Test Slide"

    # Verify shapes extraction
    assert len(slide_data["shapes"]) == 2
    assert slide_data["shapes"][0]["shape_id"] == 1
    assert slide_data["shapes"][1]["shape_id"] == 2


# Test PowerPointShapeExtractor
def test_powerpoint_shape_extractor():
    # Create a mock presentation
    mock_ppt = MagicMock(spec=Presentation)
    mock_ppt.slide_width = Length(800)
    mock_ppt.slide_height = Length(600)
    mock_ppt.slides = [MagicMock(spec=Slide), MagicMock(spec=Slide)]

    # Set slide IDs for testing
    mock_ppt.slides[0].slide_id = 1
    mock_ppt.slides[0].name = "Slide 1"
    mock_ppt.slides[0].shapes = [MagicMock(shape_id=1, shape_type="AUTO_SHAPE")]

    mock_ppt.slides[1].slide_id = 2
    mock_ppt.slides[1].name = "Slide 2"
    mock_ppt.slides[1].shapes = [MagicMock(shape_id=2, shape_type="TEXT_BOX")]

    ppt_extractor = PowerPointShapeExtractor(mock_ppt, measurement_unit="pt")

    # Extract PowerPoint data
    ppt_data = ppt_extractor.extract_ppt()

    # Verify PowerPoint metadata
    assert ppt_data["slide_width"] == Length(800).pt
    assert ppt_data["slide_height"] == Length(600).pt

    # Verify slides extraction
    assert len(ppt_data["slides"]) == 2
    assert ppt_data["slides"][0]["slide_id"] == 1
    assert ppt_data["slides"][1]["slide_id"] == 2
    assert len(ppt_data["slides"][0]["shapes"]) == 1
    assert len(ppt_data["slides"][1]["shapes"]) == 1


# Run the tests
if __name__ == "__main__":
    pytest.main()
