import os
import sys

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Inches

sys.path.append(
    os.path.abspath(os.path.join(os.path.dirname(__file__), "../src/pptlayout/input"))
)

from shape_wrappers import BaseShapeWrapper  # noqa: E402
from shape_wrappers import PictureWrapper  # noqa: E402
from shape_wrappers import TableWrapper  # noqa: E402

PPTX_PATH = os.path.join(os.path.dirname(__file__), "./data/test.pptx")


@pytest.fixture
def load_presentation():
    """Fixture to load the presentation."""
    ppt = Presentation(PPTX_PATH)
    return ppt


def test_base_shape_wrapper():
    # Create a sample presentation
    ppt = Presentation()

    # Add a slide with a title and content layout
    slide_layout = ppt.slide_layouts[5]  # Choosing a blank slide layout
    slide = ppt.slides.add_slide(slide_layout)

    # Add a shape to the slide (rectangle, for example)
    left = Inches(1)
    top = Inches(1)
    width = Inches(2)
    height = Inches(1)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.text = "Test Shape"

    # Add another shape (e.g., ellipse)
    ellipse = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(4), Inches(2), Inches(3), Inches(2)
    )

    # Create BaseShapeWrapper instances for each shape
    shape_wrapper = BaseShapeWrapper(shape)
    ellipse_wrapper = BaseShapeWrapper(ellipse)

    # Print descriptions for both shapes
    print("Testing Shape Wrapper:")
    print(shape_wrapper)

    print("Testing Ellipse Wrapper:")
    print(ellipse_wrapper)

    # Assertions for testing
    assert "AUTO_SHAPE" in shape_wrapper.shape_description
    assert "Position: left=914400, top=914400" in shape_wrapper.position_info
    assert "Size: width=1828800, height=914400" in shape_wrapper.size_info

    assert "AUTO_SHAPE" in ellipse_wrapper.shape_description
    assert ellipse_wrapper.size_info == f"Size: width={Inches(3)}, height={Inches(2)}\n"

    print("test_base_shape_wrapper passed!")


def test_picture_wrapper(load_presentation):
    """Test the PictureWrapper with pictures from the third slide."""
    # Get the third slide (index 2 since it's zero-based)
    slide = load_presentation.slides[2]

    picture_wrappers = []

    # Loop through shapes in the third slide and create PictureWrapper instances for pictures
    for shape in slide.shapes:
        print(str(shape.shape_type))
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_wrapper = PictureWrapper(shape)
            picture_wrappers.append(picture_wrapper)

    # Assertions for testing
    assert len(picture_wrappers) > 0, "No pictures found on this slide."

    for idx, wrapper in enumerate(picture_wrappers):
        print(f"Testing Picture Wrapper {idx + 1}:")
        print(wrapper)

        # Check if the description contains "Picture"
        assert "Picture" in wrapper.shape_description
        # Check if rotation info is available
        assert isinstance(wrapper.rotation_info, str)

    print("All picture wrapper tests passed!")


def test_table_wrapper(load_presentation):
    """Test the TableWrapper with tables from the fourth slide."""
    # Get the fourth slide (index 3 since it's zero-based)
    slide = load_presentation.slides[9]

    table_wrappers = []

    # Loop through shapes in the fourth slide and create TableWrapper instances for tables
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_wrapper = TableWrapper(shape)
            table_wrappers.append(table_wrapper)

    # Assertions for testing
    assert len(table_wrappers) > 0, "No tables found on this slide."

    for idx, wrapper in enumerate(table_wrappers):
        print(f"Testing Table Wrapper {idx + 1}:")
        print(wrapper)

        # Check if the description contains "Table"
        assert "Table" in wrapper.shape_description

    print("All table wrapper tests passed!")
