from pptx import Presentation


class PowerPointToHTMLConverter:
    def __init__(self, ppt_file_path: str):
        self.ppt_file_path = ppt_file_path
        self.ppt = Presentation(ppt_file_path)
        self.html = ""

    def convert(self):
        for slide in self.ppt.slides:
            self.html += "<div>"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    self.html += f"<p>{shape.text}</p>"
            self.html += "</div>"
        return self.html
