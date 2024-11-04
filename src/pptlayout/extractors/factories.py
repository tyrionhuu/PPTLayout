from typing import TypeAlias, Union

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture
from pptx.shapes.placeholder import BasePlaceholder

from .shape_extractors import (
    BaseAutoShapeExtractor,
    BaseShapeExtractor,
    ConnectorExtractor,
    FreeformExtractor,
    GraphicFrameExtractor,
    GroupShapeExtractor,
    MovieExtractor,
    PictureExtractor,
    PlaceholderExtractor,
)

ShapeExtractor: TypeAlias = Union[
    BaseShapeExtractor,
    BaseAutoShapeExtractor,
    ConnectorExtractor,
    FreeformExtractor,
    GroupShapeExtractor,
    PictureExtractor,
    MovieExtractor,
    PlaceholderExtractor,
]

Shape: TypeAlias = Union[
    BaseShape,
    AutoShape,
    Connector,
    GraphicFrame,
    GroupShape,
    Picture,
    Movie,
    BasePlaceholder,
]

SHAPE_EXTRACTOR_MAP = {
    # Auto Shape
    MSO_SHAPE_TYPE.AUTO_SHAPE: BaseAutoShapeExtractor,
    MSO_SHAPE_TYPE.TEXT_BOX: BaseAutoShapeExtractor,
    MSO_SHAPE_TYPE.FREEFORM: FreeformExtractor,
    MSO_SHAPE_TYPE.PLACEHOLDER: PlaceholderExtractor,
    # Graphic Frame
    MSO_SHAPE_TYPE.CHART: GraphicFrameExtractor,
    MSO_SHAPE_TYPE.TABLE: GraphicFrameExtractor,
    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: GraphicFrameExtractor,
    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: GraphicFrameExtractor,
    # Picture
    MSO_SHAPE_TYPE.PICTURE: PictureExtractor,
    MSO_SHAPE_TYPE.MEDIA: MovieExtractor,
    # Connector
    MSO_SHAPE_TYPE.LINE: ConnectorExtractor,
    # Group Shape
    MSO_SHAPE_TYPE.GROUP: GroupShapeExtractor,
}

DEFAULT_EXTRACTOR = BaseShapeExtractor


def shape_extractor_factory(
    shape: Shape, measurement_unit: str = "pt"
) -> ShapeExtractor:
    """Factory function to create a shape extractor based on the shape type."""
    shape_type = shape.shape_type
    extractor = SHAPE_EXTRACTOR_MAP.get(shape_type, DEFAULT_EXTRACTOR)
    return extractor(shape, measurement_unit)
