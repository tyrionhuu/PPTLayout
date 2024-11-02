from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.util import Length


class BaseShapeHandler:
    def __init__(self, shape: BaseShape):
        self._shape = shape
        self._shape_type = shape.shape_type
        self._shape_id = shape.shape_id
        self._name = shape.name
        self._height = shape.height
        self._width = shape.width
        self._left = shape.left
        self._top = shape.top
        self._rotation = shape.rotation
        self._has_chart = shape.has_chart
        self._has_text_frame = shape.has_text_frame
        self._has_table = shape.has_table
        self._is_placeholder = shape.is_placeholder
        self._placeholder_format = shape.placeholder_format

    @property
    def shape_type(self) -> MSO_SHAPE_TYPE:
        return self._shape_type

    @property
    def shape_id(self) -> int:
        return self._shape_id

    @property
    def name(self) -> str:
        return self._name

    @property
    def height(self) -> Length:
        return self._height

    @property
    def width(self) -> Length:
        return self._width

    @property
    def left(self) -> Length:
        return self._left

    @property
    def top(self) -> Length:
        return self._top

    @property
    def rotation(self) -> float:
        return self._rotation

    @property
    def has_chart(self) -> bool:
        return self._has_chart

    @property
    def has_text_frame(self) -> bool:
        return self._has_text_frame

    @property
    def has_table(self) -> bool:
        return self._has_table

    @property
    def is_placeholder(self) -> bool:
        return self._is_placeholder

    @property
    def placeholder_format(self):
        return self._placeholder_format
