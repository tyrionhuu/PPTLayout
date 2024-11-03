from typing import Union

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.util import Length


class BaseShapeExtractor:
    def __init__(self, shape: BaseShape, measurement_unit: str = "pt"):
        self._shape = shape
        self.measurement_unit = measurement_unit

    def _extract_shape_type(self) -> str:
        shape_type = self._shape.shape_type
        # Check if the shape type is a valid MSO_SHAPE_TYPE enum member
        if isinstance(shape_type, MSO_SHAPE_TYPE):
            return shape_type.name  # Returns the name of the enum member
        return str(shape_type)  # Fallback in case it's not in the enum

    def _unit_conversion(self, value: Length, unit: str) -> Union[int, float]:
        if unit == "cm":
            return value.cm
        elif unit == "inches" or unit == "in" or unit == "inch":
            return value.inches
        elif unit == "pt":
            return value.pt
        elif unit == "emu":
            return value.emu
        else:
            raise ValueError(f"Invalid measurement unit: {unit}")

    def _extract_height(self) -> Union[int, float]:
        return self._unit_conversion(self._shape.height, self.measurement_unit)

    def _extract_width(self) -> Union[int, float]:
        return self._unit_conversion(self._shape.width, self.measurement_unit)

    def _extract_left(self) -> Union[int, float]:
        return self._unit_conversion(self._shape.left, self.measurement_unit)

    def _extract_top(self) -> Union[int, float]:
        return self._unit_conversion(self._shape.top, self.measurement_unit)

    def set_measurement_unit(self, unit: str) -> None:
        self.measurement_unit = unit

    def extract_shape(self) -> dict:
        return {
            "name": self._shape.name,
            "shape_id": self._shape.shape_id,
            "shape_type": self._extract_shape_type(),
            "measurement_unit": self.measurement_unit,
            "height": self._extract_height(),
            "width": self._extract_width(),
            "left": self._extract_left(),
            "top": self._extract_top(),
        }


class AutoShapeExtractor(BaseShapeExtractor):
    def __init__(self, shape: AutoShape, measurement_unit: str = "pt"):
        super().__init__(shape, measurement_unit)

    def _extract_auto_shape_type(self) -> str:
        auto_shape_type = self._shape.auto_shape_type  # type: ignore[attr-defined]
        # Check if the shape type is a valid MSO_AUTO_SHAPE_TYPE enum member
        if isinstance(auto_shape_type, MSO_AUTO_SHAPE_TYPE):
            return auto_shape_type.name
        raise AttributeError("Shape is not an auto shape")

    def _extract_text(self) -> str:
        if self._shape.has_text_frame:
            return self._shape.text  # type: ignore[attr-defined]
        raise AttributeError("Shape does not have a text frame")

    def extract_shape(self) -> dict:
        shape_data = super().extract_shape()
        shape_data["auto_shape_type"] = self._extract_auto_shape_type()
        if self._shape.has_text_frame:
            shape_data["text"] = self._extract_text()
        return shape_data


class ConnectorExtractor(BaseShapeExtractor):
    def __init__(self, shape: Connector, measurement_unit: str = "pt"):
        super().__init__(shape, measurement_unit)

    def _extract_begin_x(self) -> Union[int, float]:
        return self._unit_conversion(self._shape.begin_x, self.measurement_unit)  # type: ignore[attr-defined]

    def _extract_begin_y(self) -> Union[int, float]:
        return self._unit_conversion(self._shape.begin_y, self.measurement_unit)  # type: ignore[attr-defined]

    def _extract_end_x(self) -> Union[int, float]:
        return self._unit_conversion(self._shape.end_x, self.measurement_unit)  # type: ignore[attr-defined]

    def _extract_end_y(self) -> Union[int, float]:
        return self._unit_conversion(self._shape.end_y, self.measurement_unit)  # type: ignore[attr-defined]

    def extract_shape(self) -> dict:
        shape_data = super().extract_shape()
        shape_data["begin_x"] = self._extract_begin_x()
        shape_data["begin_y"] = self._extract_begin_y()
        shape_data["end_x"] = self._extract_end_x()
        shape_data["end_y"] = self._extract_end_y()
        return shape_data


class FreeformExtractor(BaseShapeExtractor):
    pass
