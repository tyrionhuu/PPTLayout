import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes import autoshape

from src.pptlayout.extractors.shape_extractors import (
    BaseAutoShapeExtractor,
    BaseShapeExtractor,
)


@pytest.fixture
def pptx_presentation():
    """Fixture to load a PowerPoint presentation."""
    PPTX_PATH = "data/pptx/ZK7FNUZ33GBBCG7CFVYS56TQCTD72CJR.pptx"
    presentation = Presentation(PPTX_PATH)
    yield presentation  # Yield the presentation for use in tests


def get_base_shape_from_pptx(presentation):
    """Function to get a base shape from a given PowerPoint presentation."""
    base_shape = None

    # Iterate through all slides in the presentation
    for slide in presentation.slides:
        # Iterate through all shapes in the slide
        for shape in slide.shapes:
            # Check if the shape is not a group shape
            if not shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                base_shape = shape
                break

    return base_shape


def get_auto_shape_from_pptx(presentation):
    """Function to get an auto shape from a given PowerPoint presentation."""
    auto_shape = None

    # Iterate through all slides in the presentation
    for slide in presentation.slides:
        # Iterate through all shapes in the slide
        for shape in slide.shapes:
            # Check if the shape is an auto shape
            if isinstance(shape, autoshape.Shape):
                auto_shape = shape
                break

    return auto_shape


def test_extract_shape(pptx_presentation):
    """Test to check the extraction of shape details from the presentation."""
    base_shape = get_base_shape_from_pptx(pptx_presentation)
    shape_extractor = BaseShapeExtractor(base_shape, measurement_unit="cm")
    shape_info = shape_extractor.extract_shape()

    # Assert that shape details are extracted
    assert shape_info is not None, "No shape details found in the presentation."
    # print(shape_info)
    # Optionally, check details of the first shape
    assert shape_info["name"] is not None, "The shape does not have a name."
    assert shape_info["shape_id"] is not None, "The shape does not have an ID."
    assert shape_info["shape_type"] is not None, "The shape does not have a type."
    assert (
        shape_info["measurement_unit"] is not None
    ), "The shape does not have a measurement unit."
    assert shape_info["height"] is not None, "The shape does not have a height."
    assert shape_info["width"] is not None, "The shape does not have a width."
    assert shape_info["left"] is not None, "The shape does not have a left position."
    assert shape_info["top"] is not None, "The shape does not have a top position."


def test_extract_auto_shapes(pptx_presentation):
    """Test to check the extraction of auto shapes from the presentation."""
    auto_shape = get_auto_shape_from_pptx(pptx_presentation)
    print(BaseAutoShapeExtractor(auto_shape, measurement_unit="inches").extract_shape())
    # Assert that auto shapes are extracted
    assert auto_shape is not None, "No auto shapes found in the presentation."

    # Optionally, check details of the first auto shape
    assert isinstance(auto_shape, autoshape.Shape), "The  shape is not an auto shape."
    assert auto_shape.name is not None, "The auto shape does not have a name."
