import os
import sys

from pptx import Presentation

sys.path.append(
    os.path.abspath(os.path.join(os.path.dirname(__file__), "../src/pptlayout/input"))
)

from utils import get_fill_color  # noqa: E402


class SlideShapeExtractor:
    """
    Extract the shapes from a slide in a presentation to a csv file, which is formatted as follows:
    Slide Number, Shape Type, Shape Index, Left Position, Top Position, Width, Height, Fill color
    """

    def __init__(self, slide, slide_number: int):
        self.slide = slide
        self.slide_number = slide_number
        self.shapes = self._extract_shapes()

    def _extract_shapes(self) -> list:
        shapes = []
        for shape in self.slide.shapes:
            shapes.append(shape)
        return shapes

    def extract_shapes_to_csv(self, output_file: str):
        with open(output_file, "w") as f:
            f.write(
                "Slide Number, Shape Type, Shape Index, Left Position, Top Position, Width, Height, Fill color\n"
            )
            for i, shape in enumerate(self.shapes):
                slide_number = self.slide_number

                try:
                    shape_type = shape.shape_type
                except AttributeError:
                    shape_type = ""

                try:
                    shape_index = shape.shape_id
                except AttributeError:
                    shape_index = ""

                try:
                    left_position = shape.left
                except AttributeError as e:
                    print(f"Error: {e}")

                try:
                    top_position = shape.top
                except AttributeError as e:
                    print(f"Error: {e}")

                try:
                    width = shape.width
                except AttributeError as e:
                    print(f"Error: {e}")

                try:
                    height = shape.height
                except AttributeError as e:
                    print(f"Error: {e}")

                fill_color = get_fill_color(shape)

                f.write(
                    f"{slide_number}, {shape_type}, {shape_index}, {left_position}, {top_position}, {width}, {height}, {fill_color}\n"
                )
        print(f"Shapes extracted to {output_file}")


def extract_shapes_from_presentation(presentation_file: str, output_file: str):
    presentation = Presentation(presentation_file)
    for i, slide in enumerate(presentation.slides):
        slide_extractor = SlideShapeExtractor(slide, i)
        slide_extractor.extract_shapes_to_csv(output_file)

    print(f"Shapes extracted from {presentation_file} to {output_file}")
