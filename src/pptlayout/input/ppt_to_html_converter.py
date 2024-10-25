from pptx import Presentation


class PowerPointToHTMLConverter:
    def __init__(self, ppt_file_path: str):
        self.ppt_file_path = ppt_file_path
        self.ppt = Presentation(ppt_file_path)
        self.htmls = [str]
        self.slides = self.ppt.slides
